APP_VERSION = "2.4"
import sys
import os
import math
import json
import shutil
import requests
import subprocess
import webbrowser
from datetime import datetime, timezone
from email.utils import parsedate_to_datetime
from PyQt6.QtWidgets import (QApplication, QGraphicsView, QGraphicsScene, QMainWindow, 
                             QGraphicsPixmapItem, QGraphicsTextItem, QVBoxLayout, 
                             QWidget, QPushButton, QHBoxLayout, QSlider, QLineEdit, 
                             QFormLayout, QColorDialog, QGraphicsPathItem, QCheckBox,
                             QGroupBox, QGraphicsColorizeEffect, QMessageBox, QListWidget,
                             QGraphicsEllipseItem, QGraphicsItem, QComboBox, QFileDialog, 
                             QGraphicsRectItem, QLabel, QGraphicsLineItem, QGraphicsPolygonItem, 
                             QButtonGroup, QListWidgetItem, QTabWidget, QInputDialog, QDialog, QMenu, QGridLayout, QScrollArea, QSizePolicy)
from PyQt6.QtGui import (QPen, QBrush, QColor, QPixmap, QFont, QPainter, 
                         QPainterPath, QImage, QPolygonF, QPainterPathStroker, QIcon, QTransform, QAction,
                         QPdfWriter, QPageSize, QPageLayout, QKeySequence, QShortcut)
from PyQt6.QtCore import Qt, QPointF, QRectF, QTimer, QSize

try:
    from pptx import Presentation
    from pptx.util import Inches
    import cv2
    import numpy as np
except ImportError:
    print("Kütüphaneler eksik. Terminale: pip install python-pptx opencv-python numpy Pillow")

import ctypes

os.environ["QT_ENABLE_HIGHDPI_SCALING"] = "1"
os.environ["QT_AUTO_SCREEN_SCALE_FACTOR"] = "1"
try:
    ctypes.windll.shcore.SetProcessDpiAwareness(1)
except:
    pass

def resource_path(relative_path):
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")
    return os.path.join(base_path, relative_path)

if getattr(sys, 'frozen', False):
    APP_DIR = os.path.dirname(sys.executable)
else:
    APP_DIR = os.path.dirname(os.path.abspath(__file__))

AUTH_FILE = os.path.join(os.path.expanduser("~"), ".footyscopefg_auth.json")

def get_asset(name):
    base_name = name.replace("..png", ".png")
    try:
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")
        
    search_paths = [
        os.path.join(base_path, base_name),
        os.path.join(base_path, "Futbol_Analiz", base_name),
        os.path.join(APP_DIR, base_name),
        os.path.join(APP_DIR, "Futbol_Analiz", base_name),
        os.path.join(os.getcwd(), base_name),
        os.path.join(os.getcwd(), "Futbol_Analiz", base_name),
        resource_path(name.replace(".png", "..png")) 
    ]
    for p in search_paths:
        if os.path.exists(p): return p
    return resource_path(base_name) 

PLAYER_ASSET = get_asset("player_asset.png")
GK_ASSET = get_asset("gk_asset.png")

MEMORY_DIR = os.path.join(APP_DIR, "saved_players")
TACTICS_DIR = os.path.join(APP_DIR, "saved_tactics")
VIDEO_DIR = os.path.join(APP_DIR, "saved_videos")
os.makedirs(MEMORY_DIR, exist_ok=True)
os.makedirs(TACTICS_DIR, exist_ok=True)
os.makedirs(VIDEO_DIR, exist_ok=True)

IMAGE_CACHE = {}

LANG = {
    "Türkçe": {
        "title": "FootyscopeFG ⎪ Taktik Maç Kartı Analiz Paneli",
        "pitch_group": "🏟️ Saha & Zemin", "dir": "Yön:", "surface": "Zemin:",
        "pitch_types": ["Yatay Tam Saha", "Dikey Tam Saha", "Yatay Yarım Saha", "Dikey Yarım Saha"],
        "grass_types": ["Çizgili", "Damalı", "Düz"],
        "home": "🏠 Ev Sahibi", "away": "✈️ Deplasman", "tactic": "Taktik:", "select_tactic": "Taktik Seç...",
        "color1": "Renk 1", "color2": "Renk 2", "apply": "⚡ Sahaya Uygula",
        "management": "⚙️ Yönetim", "save": "💾 Kaydet", "pptx": "📊 PPTX Çıktı", "pdf": "📄 PDF Çıktı",
        "record": "🔴 Video Kayıt", "stop_record": "⏹️ Durdur", "logout": "🚪 Çıkış Yap",
        "tools": "Araçlar", "tool_select": "👆 Seç / Taşı", "tool_h_man": "🧑 Ev (Manuel)", "tool_a_man": "🧑 Dep (Manuel)",
        "tool_arr_s": "➡️ Düz Ok", "tool_arr_c": "↪️ Kavisli Ok", "tool_rect": "⬛ Kare", "tool_circ": "⏺ Daire",
        "tool_text": "📝 Metin", "tool_ball": "⚽ Top Ekle",
        "balls": ["Adidas Klasik", "Şampiyonlar Ligi", "Dünya Kupası", "Premier Lig", "Nike Pro"],
        "edit_obj": "Nesne Düzenleme", "name": "İsim:", "color": "Renk:", "size": "Boyut:", "angle": "Açı:",
        "mirror": "Ayna:", "mir_h": "↔️ Yatay", "mir_v": "↕️ Dikey", "curve": "Kavis:",
        "other": "🎨 Diğer", "replace_img": "🖼️ Figür Görseli Değiştir", "delete_sel": "🗑️ Seçileni Sil",
        "tab_img": "🖼️ Görseller", "tab_tac": "📋 Taktikler", "tab_vid": "🎥 Videolar",
        "add_pc": "💻 Bilgisayardan Ekle", "add_vid": "💻 Video Ekle", "load": "📂 Yükle",
        "ev_pre": "EV", "dep_pre": "DEP",
        "ctx_del_pc": "🗑️ Bilgisayardan Sil", "ctx_play": "▶️ Aç / Oynat", "ctx_ren": "✏️ Yeniden Adlandır",
        "ctx_del_vid": "🗑️ Sil", "ctx_del_tac": "🗑️ Taktiği Sil",
        "color_target": "Hedef:", "col_target_jersey": "👕 Forma Rengi", "col_target_text": "🔤 Yazı Rengi"
    },
    "English": {
        "title": "FootyscopeFG ⎪ Tactical Match Card Analysis Panel",
        "pitch_group": "🏟️ Pitch & Surface", "dir": "Direction:", "surface": "Surface:",
        "pitch_types": ["Horizontal Full", "Vertical Full", "Horizontal Half", "Vertical Half"],
        "grass_types": ["Striped", "Checkered", "Plain"],
        "home": "🏠 Home", "away": "✈️ Away", "tactic": "Tactic:", "select_tactic": "Select Tactic...",
        "color1": "Color 1", "color2": "Color 2", "apply": "⚡ Apply to Pitch",
        "management": "⚙️ Management", "save": "💾 Save", "pptx": "📊 PPTX Export", "pdf": "📄 PDF Export",
        "record": "🔴 Record Video", "stop_record": "⏹️ Stop", "logout": "🚪 Logout",
        "tools": "Tools", "tool_select": "👆 Select / Move", "tool_h_man": "🧑 Home (Manual)", "tool_a_man": "🧑 Away (Manual)",
        "tool_arr_s": "➡️ Straight Arrow", "tool_arr_c": "↪️ Curved Arrow", "tool_rect": "⬛ Square", "tool_circ": "⏺ Circle",
        "tool_text": "📝 Text", "tool_ball": "⚽ Add Ball",
        "balls": ["Adidas Classic", "Champions League", "World Cup", "Premier League", "Nike Pro"],
        "edit_obj": "Object Editing", "name": "Name:", "color": "Color:", "size": "Size:", "angle": "Angle:",
        "mirror": "Mirror:", "mir_h": "↔️ Horizontal", "mir_v": "↕️ Vertical", "curve": "Curve:",
        "other": "🎨 Other", "replace_img": "🖼️ Replace Figure Image", "delete_sel": "🗑️ Delete Selected",
        "tab_img": "🖼️ Images", "tab_tac": "📋 Tactics", "tab_vid": "🎥 Videos",
        "add_pc": "💻 Add from PC", "add_vid": "💻 Add Video", "load": "📂 Load",
        "ev_pre": "HM", "dep_pre": "AW",
        "ctx_del_pc": "🗑️ Delete from PC", "ctx_play": "▶️ Play", "ctx_ren": "✏️ Rename",
        "ctx_del_vid": "🗑️ Delete", "ctx_del_tac": "🗑️ Delete Tactic",
        "color_target": "Target:", "col_target_jersey": "👕 Jersey Color", "col_target_text": "🔤 Text Color"
    },
    "Español": {
        "title": "FootyscopeFG ⎪ Panel de Análisis Táctico",
        "pitch_group": "🏟️ Campo y Superficie", "dir": "Dirección:", "surface": "Césped:",
        "pitch_types": ["Horizontal Completo", "Vertical Completo", "Horizontal Medio", "Vertical Medio"],
        "grass_types": ["Rayado", "Ajedrezado", "Liso"],
        "home": "🏠 Local", "away": "✈️ Visitante", "tactic": "Táctica:", "select_tactic": "Elegir Táctica...",
        "color1": "Color 1", "color2": "Color 2", "apply": "⚡ Aplicar al Campo",
        "management": "⚙️ Gestión", "save": "💾 Guardar", "pptx": "📊 Exportar PPTX", "pdf": "📄 Exportar PDF",
        "record": "🔴 Grabar Video", "stop_record": "⏹️ Detener", "logout": "🚪 Cerrar Sesión",
        "tools": "Herramientas", "tool_select": "👆 Seleccionar / Mover", "tool_h_man": "🧑 Local (Manual)", "tool_a_man": "🧑 Visitante (Manual)",
        "tool_arr_s": "➡️ Flecha Recta", "tool_arr_c": "↪️ Flecha Curva", "tool_rect": "⬛ Cuadrado", "tool_circ": "⏺ Círculo",
        "tool_text": "📝 Texto", "tool_ball": "⚽ Añadir Balón",
        "balls": ["Adidas Clásico", "Champions League", "Copa Mundial", "Premier League", "Nike Pro"],
        "edit_obj": "Edición de Objetos", "name": "Nombre:", "color": "Color:", "size": "Tamaño:", "angle": "Ángulo:",
        "mirror": "Espejo:", "mir_h": "↔️ Horizontal", "mir_v": "↕️ Vertical", "curve": "Curva:",
        "other": "🎨 Otro", "replace_img": "🖼️ Reemplazar Figura", "delete_sel": "🗑️ Eliminar Seleccionado",
        "tab_img": "🖼️ Imágenes", "tab_tac": "📋 Tácticas", "tab_vid": "🎥 Videos",
        "add_pc": "💻 Añadir de PC", "add_vid": "💻 Añadir Video", "load": "📂 Cargar",
        "ev_pre": "LC", "dep_pre": "VS",
        "ctx_del_pc": "🗑️ Eliminar de PC", "ctx_play": "▶️ Reproducir", "ctx_ren": "✏️ Renombrar",
        "ctx_del_vid": "🗑️ Eliminar", "ctx_del_tac": "🗑️ Eliminar Táctica",
        "color_target": "Objetivo:", "col_target_jersey": "👕 Color de Forma", "col_target_text": "🔤 Color de Texto"
    },
    "Deutsch": {
        "title": "FootyscopeFG ⎪ Taktik-Analyse-Panel",
        "pitch_group": "🏟️ Spielfeld & Rasen", "dir": "Richtung:", "surface": "Rasen:",
        "pitch_types": ["Horizontal Ganz", "Vertikal Ganz", "Horizontal Halb", "Vertikal Halb"],
        "grass_types": ["Gestreift", "Kariert", "Einfach"],
        "home": "🏠 Heim", "away": "✈️ Auswärts", "tactic": "Taktik:", "select_tactic": "Taktik wählen...",
        "color1": "Farbe 1", "color2": "Farbe 2", "apply": "⚡ Anwenden",
        "management": "⚙️ Verwaltung", "save": "💾 Speichern", "pptx": "📊 PPTX Export", "pdf": "📄 PDF Export",
        "record": "🔴 Video Aufnehmen", "stop_record": "⏹️ Stopp", "logout": "🚪 Abmelden",
        "tools": "Werkzeuge", "tool_select": "👆 Wählen / Bewegen", "tool_h_man": "🧑 Heim (Manuell)", "tool_a_man": "🧑 Auswärts (Manuell)",
        "tool_arr_s": "➡️ Gerade Pfeil", "tool_arr_c": "↪️ Gebogener Pfeil", "tool_rect": "⬛ Rechteck", "tool_circ": "⏺ Kreis",
        "tool_text": "📝 Text", "tool_ball": "⚽ Ball hinzufügen",
        "balls": ["Adidas Klassisch", "Champions League", "Weltmeisterschaft", "Premier League", "Nike Pro"],
        "edit_obj": "Objektbearbeitung", "name": "Name:", "color": "Farbe:", "size": "Größe:", "angle": "Winkel:",
        "mirror": "Spiegeln:", "mir_h": "↔️ Horizontal", "mir_v": "↕️ Vertikal", "curve": "Kurve:",
        "other": "🎨 Andere", "replace_img": "🖼️ Figur ersetzen", "delete_sel": "🗑️ Auswahl löschen",
        "tab_img": "🖼️ Bilder", "tab_tac": "📋 Taktiken", "tab_vid": "🎥 Videos",
        "add_pc": "💻 Vom PC", "add_vid": "💻 Video", "load": "📂 Laden",
        "ev_pre": "HM", "dep_pre": "AW",
        "ctx_del_pc": "🗑️ Vom PC löschen", "ctx_play": "▶️ Abspielen", "ctx_ren": "✏️ Umbenennen",
        "ctx_del_vid": "🗑️ Löschen", "ctx_del_tac": "🗑️ Taktik löschen",
        "color_target": "Ziel:", "col_target_jersey": "👕 Trikotfarbe", "col_target_text": "🔤 Textfarbe"
    }
}

LUXURY_THEME_QSS = """
QMainWindow, QWidget { background-color: #2b0a11; color: #f5ebd9; font-family: 'Helvetica Neue', 'San Francisco', 'Segoe UI', Arial, sans-serif; }
QLabel { color: #d4af37; font-weight: bold; font-size: 11px; }
QGraphicsView { background-color: #1a060a; border: 2px solid #d4af37; border-radius: 8px; }
QGroupBox { border: 1px solid #b76e79; border-radius: 8px; margin-top: 10px; padding-top: 12px; background-color: #3b0e17; }
QGroupBox::title { subcontrol-origin: margin; subcontrol-position: top left; left: 10px; top: 0px; color: #d4af37; font-weight: bold; font-size: 11px; background-color: #2b0a11; padding: 2px 8px; border: 1px solid #b76e79; border-radius: 4px; }
QPushButton { background-color: #4a121d; border: 1px solid #b76e79; border-radius: 4px; padding: 4px 8px; color: #f5ebd9; font-weight: bold; min-height: 22px; font-size: 11px; }
QPushButton:hover { background-color: #5c1624; border-color: #d4af37; color: #ffffff; }
QPushButton:checked { background-color: #d4af37; color: #2b0a11; border: 1px solid #ffdf73; }
QLineEdit, QListWidget { background-color: #1a060a; border: 1px solid #b76e79; border-radius: 4px; padding: 4px; color: #f5ebd9; min-height: 22px; font-size: 11px; }
QComboBox { background-color: #4a121d; border: 1px solid #b76e79; border-radius: 4px; padding: 3px 6px; color: #f5ebd9; font-weight: bold; min-height: 22px; font-size: 11px; }
QComboBox::drop-down { border: none; }
QComboBox QAbstractItemView { background-color: #4a121d; color: #f5ebd9; selection-background-color: #d4af37; selection-color: #2b0a11; border: 1px solid #b76e79; }
QSlider::groove:horizontal { background: #1a060a; height: 6px; border-radius: 3px; border: 1px solid #b76e79; }
QSlider::handle:horizontal { background: #d4af37; width: 14px; margin: -4px 0; border-radius: 7px; border: 1px solid #fff; }
.ActionBtn { background-color: #d4af37; color: #1a060a; border: none; font-weight: 800; }
.ActionBtn:hover { background-color: #ffdf73; }
.DangerBtn { background-color: #8b0000; color: white; border: 1px solid #ff4d4d; }
.DangerBtn:hover { background-color: #b30000; }
QTabWidget::pane { border: 1px solid #d4af37; border-radius: 6px; background: #3b0e17; margin-top: 2px;}
QTabBar::tab { background: #2b0a11; color: #b76e79; padding: 5px 8px; border: 1px solid #b76e79; border-bottom: none; border-top-left-radius: 6px; border-top-right-radius: 6px; font-weight: bold; font-size: 11px; margin-right: 2px; }
QTabBar::tab:selected { background: #3b0e17; color: #d4af37; border-bottom: 2px solid #3b0e17; border-color: #d4af37; border-bottom-color: #3b0e17; }
QScrollArea { border: none; background-color: transparent; }
"""

class CustomRect(QGraphicsRectItem):
    def __init__(self, rect):
        super().__init__(rect)
        self.setPen(QPen(Qt.GlobalColor.white, 3)); self.setBrush(QBrush(QColor(255, 255, 255, 1))) 
        self.setFlags(QGraphicsItem.GraphicsItemFlag.ItemIsMovable | QGraphicsItem.GraphicsItemFlag.ItemIsSelectable)

class CustomEllipse(QGraphicsEllipseItem):
    def __init__(self, rect):
        super().__init__(rect)
        self.setPen(QPen(Qt.GlobalColor.white, 3)); self.setBrush(QBrush(QColor(255, 255, 255, 1)))
        self.setFlags(QGraphicsItem.GraphicsItemFlag.ItemIsMovable | QGraphicsItem.GraphicsItemFlag.ItemIsSelectable)

class Football(QGraphicsEllipseItem):
    def __init__(self, x, y, b_idx=0):
        super().__init__(0, 0, 20, 20)
        self.setPos(x, y); self.b_idx = b_idx; self.setTransformOriginPoint(10, 10)
        self.setFlags(QGraphicsItem.GraphicsItemFlag.ItemIsMovable | QGraphicsItem.GraphicsItemFlag.ItemIsSelectable | QGraphicsItem.GraphicsItemFlag.ItemSendsGeometryChanges)
        self.setPen(QPen(Qt.GlobalColor.black, 1.5))
    def paint(self, painter, option, widget):
        painter.setRenderHint(QPainter.RenderHint.Antialiasing)
        painter.setBrush(QBrush(Qt.GlobalColor.white)); painter.drawEllipse(0, 0, 20, 20)
        painter.setPen(QPen(Qt.GlobalColor.black, 1))
        if self.b_idx == 1:
            painter.setBrush(QBrush(QColor(0, 50, 120)))
            for i in range(8):
                angle = i * 45; path = QPainterPath()
                path.moveTo(10 + 6*math.cos(math.radians(angle)), 10 + 6*math.sin(math.radians(angle)))
                path.lineTo(10 + 3*math.cos(math.radians(angle+20)), 10 + 3*math.sin(math.radians(angle+20)))
                path.lineTo(10 + 3*math.cos(math.radians(angle-20)), 10 + 3*math.sin(math.radians(angle-20)))
                painter.drawPath(path)
        elif self.b_idx == 2:
            colors = [QColor(255, 50, 50), QColor(50, 50, 255), QColor(255, 200, 0)]
            for i in range(3):
                painter.setPen(QPen(colors[i], 2)); painter.drawArc(2, 2, 16, 16, (i*120)*16, 60*16)
            painter.setPen(QPen(Qt.GlobalColor.black, 1)); painter.drawEllipse(6, 6, 8, 8)
        elif self.b_idx == 3:
            painter.setPen(QPen(QColor(120, 0, 150), 2)); painter.drawEllipse(3, 3, 14, 14); painter.drawEllipse(7, 7, 6, 6)
        elif self.b_idx == 0:
            painter.setBrush(QBrush(Qt.GlobalColor.black))
            painter.drawRect(8, 2, 4, 4); painter.drawRect(2, 8, 4, 4); painter.drawRect(14, 8, 4, 4); painter.drawRect(8, 14, 4, 4)
        else:
            painter.setPen(QPen(QColor(255, 100, 0), 2)); painter.drawLine(2, 10, 18, 10); painter.drawArc(4, 4, 12, 12, 0, 180*16)

class EditableText(QGraphicsTextItem):
    def __init__(self, text="Yeni Metin", x=400, y=300):
        super().__init__(text)
        self.setPos(x, y); self.setFont(QFont("Arial", 16, QFont.Weight.Bold)); self.setDefaultTextColor(Qt.GlobalColor.white)
        self.setFlags(QGraphicsItem.GraphicsItemFlag.ItemIsMovable | QGraphicsItem.GraphicsItemFlag.ItemIsSelectable | QGraphicsItem.GraphicsItemFlag.ItemIsFocusable)
        self.setTextInteractionFlags(Qt.TextInteractionFlag.NoTextInteraction) 
    def mouseDoubleClickEvent(self, event):
        self.setTextInteractionFlags(Qt.TextInteractionFlag.TextEditorInteraction)
        self.setFocus()
        super().mouseDoubleClickEvent(event)
    def focusOutEvent(self, event):
        self.setTextInteractionFlags(Qt.TextInteractionFlag.NoTextInteraction)
        super().focusOutEvent(event)
        if not self.toPlainText().strip() and self.scene():
            self.scene().removeItem(self)

class Arrow(QGraphicsPathItem):
    def __init__(self, start_pos, end_pos, is_curved=False, color=Qt.GlobalColor.yellow):
        super().__init__()
        self.start_p, self.end_p = start_pos, end_pos
        self.is_curved, self.curve_offset = is_curved, 0.3
        self.color = QColor(color); self.head_item = QGraphicsPolygonItem(self)
        self.setFlags(QGraphicsItem.GraphicsItemFlag.ItemIsSelectable | QGraphicsItem.GraphicsItemFlag.ItemIsMovable)
        self.update_arrow(start_pos, end_pos)
    def update_arrow(self, start_pos, end_pos):
        self.start_p, self.end_p = start_pos, end_pos
        path = QPainterPath(start_pos)
        if self.is_curved:
            mid_x, mid_y = (start_pos.x()+end_pos.x())/2, (start_pos.y()+end_pos.y())/2
            dx, dy = end_pos.x()-start_pos.x(), end_pos.y()-start_pos.y()
            ctrl_x, ctrl_y = mid_x - dy*self.curve_offset, mid_y + dx*self.curve_offset
            path.quadTo(QPointF(ctrl_x, ctrl_y), end_pos)
            angle = math.atan2(end_pos.y()-ctrl_y, end_pos.x()-ctrl_x)
        else:
            path.lineTo(end_pos); angle = math.atan2(end_pos.y()-start_pos.y(), end_pos.x()-start_pos.x())
        self.setPath(path); self.setPen(QPen(self.color, 4, Qt.PenStyle.SolidLine, Qt.PenCapStyle.RoundCap))
        s = 14
        p1 = end_pos
        p2 = QPointF(end_pos.x()-s*math.cos(angle-math.pi/6), end_pos.y()-s*math.sin(angle-math.pi/6))
        p3 = QPointF(end_pos.x()-s*math.cos(angle+math.pi/6), end_pos.y()-s*math.sin(angle+math.pi/6))
        self.head_item.setPolygon(QPolygonF([p1, p2, p3])); self.head_item.setBrush(QBrush(self.color)); self.head_item.setPen(QPen(self.color, 1))
    def shape(self): return QPainterPathStroker(QPen(Qt.GlobalColor.black, 20)).createStroke(self.path())

class Player(QGraphicsPixmapItem):
    def __init__(self, image_path, name, color1, color2, x, y, is_gk=False, team=None, is_custom=False):
        super().__init__()
        self.is_gk, self.is_mirrored_h, self.is_mirrored_v = is_gk, False, False
        self.image_path, self.color1, self.color2 = image_path, QColor(color1), QColor(color2) if color2 else None
        self.team = team
        self.is_custom = is_custom 
        self.setTransformOriginPoint(25, 25); self.setPos(x, y)
        self.setFlags(QGraphicsItem.GraphicsItemFlag.ItemIsMovable | QGraphicsItem.GraphicsItemFlag.ItemIsSelectable | QGraphicsItem.GraphicsItemFlag.ItemSendsGeometryChanges)
        self.name_tag = QGraphicsTextItem(name, self); self.name_tag.setDefaultTextColor(Qt.GlobalColor.white)
        self.name_tag.setFont(QFont("Arial", 10, QFont.Weight.Bold)); self.name_tag.setFlag(QGraphicsItem.GraphicsItemFlag.ItemIgnoresTransformations) 
        self.load_image(image_path)
        
    def load_image(self, path):
        if not os.path.exists(path) and not self.is_custom: path = PLAYER_ASSET
        self.image_path = path
        cache_key = f"{path}_{self.color1.name()}_{self.color2.name() if self.color2 else 'none'}_{self.is_custom}"
        
        if cache_key in IMAGE_CACHE:
            self.original_pixmap = IMAGE_CACHE[cache_key]
        else:
            img = QImage(path).convertToFormat(QImage.Format.Format_ARGB32)
            if img.isNull() and not self.is_custom:
                img = QImage(50, 50, QImage.Format.Format_ARGB32); img.fill(Qt.GlobalColor.transparent)
                p = QPainter(img); p.setRenderHint(QPainter.RenderHint.Antialiasing); path_obj = QPainterPath()
                path_obj.moveTo(15, 10); path_obj.lineTo(35, 10); path_obj.lineTo(45, 20); path_obj.lineTo(40, 25)
                path_obj.lineTo(35, 22); path_obj.lineTo(35, 45); path_obj.lineTo(15, 45); path_obj.lineTo(15, 22)
                path_obj.lineTo(10, 25); path_obj.lineTo(5, 20); path_obj.closeSubpath()
                p.setBrush(QBrush(Qt.GlobalColor.white)); p.setPen(QPen(Qt.GlobalColor.black, 1)); p.drawPath(path_obj); p.end()

            if not self.is_custom:
                w, h = img.width(), img.height()
                for y in range(h):
                    for x in range(w):
                        c = img.pixelColor(x, y)
                        if c.alpha() > 20 and not (c.red() > 230 and c.green() > 230 and c.blue() > 230):
                            if self.color2:
                                if (x // 10) % 2 == 0: c = QColor(self.color1.red(), self.color1.green(), self.color1.blue(), c.alpha())
                                else: c = QColor(self.color2.red(), self.color2.green(), self.color2.blue(), c.alpha())
                            else: c = QColor(self.color1.red(), self.color1.green(), self.color1.blue(), c.alpha())
                            img.setPixelColor(x, y, c)
                        elif c.red() > 230 and c.green() > 230 and c.blue() > 230:
                            img.setPixelColor(x, y, QColor(255, 255, 255, 0)) 
            self.original_pixmap = QPixmap.fromImage(img).scaled(50, 50, Qt.AspectRatioMode.KeepAspectRatio, Qt.TransformationMode.SmoothTransformation)
            IMAGE_CACHE[cache_key] = self.original_pixmap
        self.set_mirrored(self.is_mirrored_h, self.is_mirrored_v); self.update_name_pos()

    def update_name_pos(self):
        n_r = self.name_tag.boundingRect(); self.name_tag.setPos((50 - n_r.width()) / 2, 50)
    def set_mirrored(self, h, v):
        self.is_mirrored_h, self.is_mirrored_v = h, v
        img = self.original_pixmap.toImage()
        if h or v: img = img.mirrored(horizontal=h, vertical=v)
        self.setPixmap(QPixmap.fromImage(img)); self.update_name_pos()

class TacticsBoard(QGraphicsView):
    def __init__(self, main_window):
        super().__init__()
        self.mw = main_window 
        self.scene = QGraphicsScene(self); self.setScene(self.scene); self.setRenderHint(QPainter.RenderHint.Antialiasing)
        self.setCacheMode(QGraphicsView.CacheModeFlag.CacheBackground)
        self.setViewportUpdateMode(QGraphicsView.ViewportUpdateMode.BoundingRectViewportUpdate)
        self.setOptimizationFlag(QGraphicsView.OptimizationFlag.DontAdjustForAntialiasing, True)
        self.pitch_items, self.cam_rotation, self.cam_3d = [], 0, False
        self.current_mode, self.temp_item, self.start_pos = 0, None, None
        self.p_w, self.p_h = 850, 550
        self.draw_pitch(0, 0) 

    def resizeEvent(self, event):
        super().resizeEvent(event)
        self.fitInView(self.scene.sceneRect(), Qt.AspectRatioMode.KeepAspectRatio)

    def update_camera(self):
        self.resetTransform(); self.rotate(self.cam_rotation)
        if self.cam_3d: self.scale(1.0, 0.65)
        self.fitInView(self.scene.sceneRect(), Qt.AspectRatioMode.KeepAspectRatio)

    def create_arc(self, x, y, w, h, start_angle, sweep_angle, pen):
        path = QPainterPath(); path.arcMoveTo(x, y, w, h, start_angle); path.arcTo(x, y, w, h, start_angle, sweep_angle)
        return self.scene.addPath(path, pen)

    def draw_pitch(self, p_idx=0, g_idx=0):
        for item in self.pitch_items: self.scene.removeItem(item)
        self.pitch_items.clear()
        
        w, h = 850, 550
        if p_idx == 0: w, h = 850, 550
        elif p_idx == 1: w, h = 550, 850
        elif p_idx == 2: w, h = 425, 550
        elif p_idx == 3: w, h = 550, 425
        
        self.p_w, self.p_h = w, h
        self.scene.setSceneRect(-150, -150, w+300, h+300)
        self.scene.setBackgroundBrush(QBrush(QColor(26, 6, 10))) 
        
        for i, color in enumerate([QColor(59, 14, 23), QColor(43, 10, 17), QColor(26, 6, 10)]):
            off = 100 - (i*30); path = QPainterPath(); path.addRoundedRect(QRectF(-off, -off, w+(off*2), h+(off*2)), 30, 30)
            p = self.scene.addPath(path, QPen(QColor(20, 20, 20), 2), QBrush(color)); p.setZValue(-40); self.pitch_items.append(p)
        
        base_color = QColor(35, 134, 54); alt_color = QColor(25, 110, 40)
        base_rect = self.scene.addRect(0, 0, w, h, QPen(Qt.GlobalColor.white, 2), QBrush(base_color)); base_rect.setZValue(-30); self.pitch_items.append(base_rect)
        
        step = 50
        if g_idx == 1: 
            for y in range(0, int(h), step):
                for x in range(0, int(w), step):
                    if ((x//step) + (y//step)) % 2 != 0: r = self.scene.addRect(x, y, step, step, QPen(Qt.PenStyle.NoPen), QBrush(alt_color)); r.setZValue(-29); self.pitch_items.append(r)
        elif g_idx == 0: 
            if p_idx in [0, 2]: 
                for x in range(0, int(w), step):
                    if (x//step) % 2 != 0: r = self.scene.addRect(x, 0, step, h, QPen(Qt.PenStyle.NoPen), QBrush(alt_color)); r.setZValue(-29); self.pitch_items.append(r)
            else: 
                for y in range(0, int(h), step):
                    if (y//step) % 2 != 0: r = self.scene.addRect(0, y, w, step, QPen(Qt.PenStyle.NoPen), QBrush(alt_color)); r.setZValue(-29); self.pitch_items.append(r)
        
        pen_w = QPen(Qt.GlobalColor.white, 2); brush_net = QBrush(QColor(255, 255, 255, 90), Qt.BrushStyle.CrossPattern)
        self.pitch_items.append(self.scene.addRect(0, 0, w, h, pen_w))
        
        if p_idx in [0, 2]:
            self.pitch_items.extend([
                self.scene.addRect(0, h/2-150, 130, 300, pen_w), self.scene.addRect(0, h/2-60, 40, 120, pen_w),
                self.scene.addRect(-35, h/2-40, 35, 80, pen_w, brush_net),
                self.scene.addEllipse(90-3, h/2-3, 6, 6, pen_w, QBrush(Qt.GlobalColor.white)),
                self.create_arc(90-73, h/2-73, 146, 146, -57, 114, pen_w)
            ])
            if p_idx == 0: 
                mid_x = w/2
                self.pitch_items.extend([self.scene.addLine(mid_x, 0, mid_x, h, pen_w), self.scene.addEllipse(mid_x-60, h/2-60, 120, 120, pen_w)])
                self.pitch_items.append(self.scene.addEllipse(mid_x-3, h/2-3, 6, 6, pen_w, QBrush(Qt.GlobalColor.white)))
                self.pitch_items.extend([
                    self.scene.addRect(w-130, h/2-150, 130, 300, pen_w), self.scene.addRect(w-40, h/2-60, 40, 120, pen_w),
                    self.scene.addRect(w, h/2-40, 35, 80, pen_w, brush_net),
                    self.scene.addEllipse(w-90-3, h/2-3, 6, 6, pen_w, QBrush(Qt.GlobalColor.white)),
                    self.create_arc(w-90-73, h/2-73, 146, 146, 123, 114, pen_w)
                ])
            elif p_idx == 2:
                self.pitch_items.extend([self.scene.addLine(w, 0, w, h, pen_w)])
                self.pitch_items.append(self.create_arc(w-60, h/2-60, 120, 120, 90, 180, pen_w))

        else: 
            self.pitch_items.extend([
                self.scene.addRect(w/2-150, 0, 300, 130, pen_w), self.scene.addRect(w/2-60, 0, 120, 40, pen_w),
                self.scene.addRect(w/2-40, -35, 80, 35, pen_w, brush_net),
                self.scene.addEllipse(w/2-3, 90-3, 6, 6, pen_w, QBrush(Qt.GlobalColor.white)),
                self.create_arc(w/2-73, 90-73, 146, 146, -147, 114, pen_w)
            ])
            if p_idx == 1:
                mid_y = h/2
                self.pitch_items.extend([self.scene.addLine(0, mid_y, w, mid_y, pen_w), self.scene.addEllipse(w/2-60, mid_y-60, 120, 120, pen_w)])
                self.pitch_items.append(self.scene.addEllipse(w/2-3, mid_y-3, 6, 6, pen_w, QBrush(Qt.GlobalColor.white)))
                self.pitch_items.extend([
                    self.scene.addRect(w/2-150, h-130, 300, 130, pen_w), self.scene.addRect(w/2-60, h-40, 120, 40, pen_w),
                    self.scene.addRect(w/2-40, h, 80, 35, pen_w, brush_net),
                    self.scene.addEllipse(w/2-3, h-90-3, 6, 6, pen_w, QBrush(Qt.GlobalColor.white)),
                    self.create_arc(w/2-73, h-90-73, 146, 146, 33, 114, pen_w)
                ])
            elif p_idx == 3: 
                self.pitch_items.extend([self.scene.addLine(0, h, w, h, pen_w)])
                self.pitch_items.append(self.create_arc(w/2-60, h-60, 120, 120, 0, 180, pen_w))

        corners = [(0,0, 270), (0,h, 0)]
        if p_idx in [0, 1]: corners.extend([(w,0, 180), (w,h, 90)])
        elif p_idx == 3: corners.extend([(w,0, 180)])
        for cx, cy, sa in corners:
            self.pitch_items.append(self.create_arc(cx-15, cy-15, 30, 30, sa, 90, pen_w))
            self.pitch_items.append(self.scene.addLine(cx, cy, cx-10, cy-15, QPen(Qt.GlobalColor.yellow, 3)))
            self.pitch_items.append(self.scene.addPolygon(QPolygonF([QPointF(cx-10, cy-15), QPointF(cx-10, cy-5), QPointF(cx-20, cy-10)]), QPen(Qt.GlobalColor.red), QBrush(Qt.GlobalColor.red)))

    def keyPressEvent(self, event):
        focus_item = self.scene.focusItem()
        if isinstance(focus_item, EditableText) and focus_item.textInteractionFlags() == Qt.TextInteractionFlag.TextEditorInteraction:
            super().keyPressEvent(event); return
            
        if event.key() in (Qt.Key.Key_Delete, Qt.Key.Key_Backspace):
            for item in self.scene.selectedItems(): self.scene.removeItem(item)
            self.mw.push_state()
        elif event.key() == Qt.Key.Key_Escape: self.current_mode = 0; self.scene.clearSelection()
        elif event.key() == Qt.Key.Key_Left: self.cam_rotation -= 90; self.update_camera()
        elif event.key() == Qt.Key.Key_Right: self.cam_rotation += 90; self.update_camera()
        elif event.key() == Qt.Key.Key_Up: self.scale(1.1, 1.1)
        elif event.key() == Qt.Key.Key_Down: self.scale(0.9, 0.9)
        elif event.key() == Qt.Key.Key_W: self.cam_3d = True; self.update_camera()
        elif event.key() == Qt.Key.Key_S: self.cam_3d = False; self.update_camera()
        else: super().keyPressEvent(event)

    def mousePressEvent(self, event):
        if self.current_mode in [1, 2, 4, 5]:
            self.start_pos = self.mapToScene(event.pos())
            if self.current_mode == 1: self.temp_item = Arrow(self.start_pos, self.start_pos, False)
            elif self.current_mode == 2: self.temp_item = Arrow(self.start_pos, self.start_pos, True)
            elif self.current_mode == 4: self.temp_item = CustomRect(QRectF(self.start_pos, self.start_pos))
            elif self.current_mode == 5: self.temp_item = CustomEllipse(QRectF(self.start_pos, self.start_pos))
            self.scene.addItem(self.temp_item)
        elif self.current_mode == 3:
            p = self.mapToScene(event.pos())
            txt_item = EditableText("Yeni Metin", p.x(), p.y())
            self.scene.addItem(txt_item); txt_item.setFocus(); self.current_mode = 0
            self.mw.push_state()
        elif self.current_mode in [7, 8]:
            self.start_pos = self.mapToScene(event.pos())
            is_h = (self.current_mode == 7)
            c1 = self.mw.b_h1.property("hex") if is_h else self.mw.b_a1.property("hex")
            c2 = self.mw.b_h2.property("hex") if is_h else self.mw.b_a2.property("hex")
            team = "home" if is_h else "away"
            pre = self.mw.tr["ev_pre"] + " " if is_h else self.mw.tr["dep_pre"] + " "
            p = Player(PLAYER_ASSET, f"{pre}X", c1, c2, self.start_pos.x(), self.start_pos.y(), False, team)
            self.scene.addItem(p)
            self.current_mode = 0; self.mw.grp.buttons()[0].setChecked(True)
            self.mw.push_state()
        else: super().mousePressEvent(event)

    def mouseMoveEvent(self, event):
        if self.temp_item and self.start_pos:
            curr = self.mapToScene(event.pos())
            if isinstance(self.temp_item, Arrow): self.temp_item.update_arrow(self.start_pos, curr)
            else: self.temp_item.setRect(QRectF(self.start_pos, curr).normalized())
        super().mouseMoveEvent(event)

    def mouseReleaseEvent(self, event):
        had_temp = self.temp_item is not None
        self.temp_item = None
        super().mouseReleaseEvent(event)
        if had_temp or self.current_mode == 0:
            self.mw.push_state()

class MainWindow(QMainWindow):
    def __init__(self, user_email=""):
        super().__init__()
        self.user_email = user_email
        self.tr = LANG["Türkçe"]
        
        self.setWindowTitle(self.tr["title"])
        self.setStyleSheet(LUXURY_THEME_QSS)
        
        self.history = []
        self.history_index = -1
        
        self.board = TacticsBoard(self); self.board.scene.selectionChanged.connect(self.on_selection_changed); self.active_item = None
        self.is_recording = False; self.record_timer = QTimer(self); self.record_timer.timeout.connect(self.record_frame); self.current_ball_type = 0
        central = QWidget(); self.setCentralWidget(central); main_layout = QVBoxLayout(central)

        top = QWidget(); top_l = QHBoxLayout(top); top_l.setContentsMargins(5, 5, 5, 5); top_l.setSpacing(8)
        
        self.g1 = QGroupBox(self.tr["pitch_group"]); l1 = QVBoxLayout(self.g1); l1.setSpacing(4); l1.setContentsMargins(5, 15, 5, 5)
        row1 = QHBoxLayout(); self.lbl_dir = QLabel(self.tr["dir"]); row1.addWidget(self.lbl_dir); self.c_pitch = QComboBox(); self.c_pitch.addItems(self.tr["pitch_types"]); self.c_pitch.currentIndexChanged.connect(self.refresh_pitch); row1.addWidget(self.c_pitch)
        row2 = QHBoxLayout(); self.lbl_surf = QLabel(self.tr["surface"]); row2.addWidget(self.lbl_surf); self.c_grass = QComboBox(); self.c_grass.addItems(self.tr["grass_types"]); self.c_grass.currentIndexChanged.connect(self.refresh_pitch); row2.addWidget(self.c_grass)
        l1.addLayout(row1); l1.addLayout(row2); top_l.addWidget(self.g1)
        
        self.g_h = QGroupBox(self.tr["home"]); l_h = QGridLayout(self.g_h); l_h.setSpacing(4); l_h.setContentsMargins(5, 15, 5, 5)
        self.c_f_h = QComboBox(); self.c_f_h.addItem(self.tr["select_tactic"]); self.c_f_h.addItems(["4-4-2", "4-3-3", "3-5-2", "4-1-4-1", "4-2-3-1", "3-4-1-2"])
        self.lbl_logo_h = QLabel(); self.lbl_logo_h.setFixedSize(24, 24); self.lbl_logo_h.setStyleSheet("border: 1px dashed #b76e79;")
        self.btn_logo_h = QPushButton("🛡️"); self.btn_logo_h.setFixedWidth(30); self.btn_logo_h.clicked.connect(lambda: self.load_team_logo("home"))
        self.b_h1 = QPushButton(self.tr["color1"]); self.b_h2 = QPushButton(self.tr["color2"])
        self.set_btn_color(self.b_h1, "#e74c3c"); self.set_btn_color(self.b_h2, "#ffffff")
        self.b_h1.clicked.connect(lambda: self.pick_team_color(self.b_h1)); self.b_h2.clicked.connect(lambda: self.pick_team_color(self.b_h2))
        row_colors_h = QHBoxLayout(); row_colors_h.setSpacing(10); row_colors_h.addWidget(self.b_h1, 1); row_colors_h.addWidget(self.b_h2, 1)
        self.b_a_h = QPushButton(self.tr["apply"]); self.b_a_h.setProperty("class", "ActionBtn"); self.b_a_h.clicked.connect(lambda: self.apply_formation("home"))
        l_h.addWidget(self.c_f_h, 0, 0, 1, 2); l_h.addWidget(self.lbl_logo_h, 0, 2); l_h.addWidget(self.btn_logo_h, 0, 3)
        l_h.addLayout(row_colors_h, 1, 0, 1, 4); l_h.addWidget(self.b_a_h, 2, 0, 1, 4); top_l.addWidget(self.g_h)

        self.g_a = QGroupBox(self.tr["away"]); l_a = QGridLayout(self.g_a); l_a.setSpacing(4); l_a.setContentsMargins(5, 15, 5, 5)
        self.c_f_a = QComboBox(); self.c_f_a.addItem(self.tr["select_tactic"]); self.c_f_a.addItems(["4-4-2", "4-3-3", "3-5-2", "4-1-4-1", "4-2-3-1", "3-4-1-2"])
        self.lbl_logo_a = QLabel(); self.lbl_logo_a.setFixedSize(24, 24); self.lbl_logo_a.setStyleSheet("border: 1px dashed #b76e79;")
        self.btn_logo_a = QPushButton("🛡️"); self.btn_logo_a.setFixedWidth(30); self.btn_logo_a.clicked.connect(lambda: self.load_team_logo("away"))
        self.b_a1 = QPushButton(self.tr["color1"]); self.b_a2 = QPushButton(self.tr["color2"])
        self.set_btn_color(self.b_a1, "#3498db"); self.set_btn_color(self.b_a2, "#2c3e50")
        self.b_a1.clicked.connect(lambda: self.pick_team_color(self.b_a1)); self.b_a2.clicked.connect(lambda: self.pick_team_color(self.b_a2))
        row_colors_a = QHBoxLayout(); row_colors_a.setSpacing(10); row_colors_a.addWidget(self.b_a1, 1); row_colors_a.addWidget(self.b_a2, 1)
        self.b_a_a = QPushButton(self.tr["apply"]); self.b_a_a.setProperty("class", "ActionBtn"); self.b_a_a.clicked.connect(lambda: self.apply_formation("away"))
        l_a.addWidget(self.c_f_a, 0, 0, 1, 2); l_a.addWidget(self.lbl_logo_a, 0, 2); l_a.addWidget(self.btn_logo_a, 0, 3)
        l_a.addLayout(row_colors_a, 1, 0, 1, 4); l_a.addWidget(self.b_a_a, 2, 0, 1, 4); top_l.addWidget(self.g_a)

        self.g_o = QGroupBox(self.tr["management"]); l_o = QGridLayout(self.g_o); l_o.setSpacing(4); l_o.setContentsMargins(5, 15, 5, 5)
        self.b_s = QPushButton(self.tr["save"]); self.b_s.clicked.connect(self.save_tactics)
        self.b_p = QPushButton(self.tr["pptx"]); self.b_p.clicked.connect(self.export_to_pptx)
        self.b_pdf = QPushButton(self.tr["pdf"]); self.b_pdf.clicked.connect(self.export_to_pdf)
        self.b_logout = QPushButton(self.tr["logout"]); self.b_logout.setProperty("class", "DangerBtn"); self.b_logout.clicked.connect(self.logout)
        
        self.b_v = QPushButton(self.tr["record"]); self.b_v.setProperty("class", "DangerBtn"); self.b_v.clicked.connect(self.toggle_recording)
        
        self.lbl_license = QLabel("⏳ Süre Çekiliyor")
        self.lbl_license.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.lbl_license.setStyleSheet("background-color: #2b0a11; color: #d4af37; border: 1px solid #b76e79; border-radius: 4px; font-weight: bold; font-size: 11px;")
        
        self.c_lang = QComboBox()
        self.c_lang.addItems(["Türkçe", "English", "Español", "Deutsch"])
        self.c_lang.currentTextChanged.connect(self.change_language)

        l_o.addWidget(self.b_s, 0, 0)
        l_o.addWidget(self.b_p, 0, 1)
        l_o.addWidget(self.b_pdf, 0, 2)
        l_o.addWidget(self.b_logout, 0, 3)
        l_o.addWidget(self.b_v, 1, 0, 1, 2)
        l_o.addWidget(self.lbl_license, 1, 2, 1, 2) 
        l_o.addWidget(self.c_lang, 2, 0, 1, 4)
        top_l.addWidget(self.g_o); main_layout.addWidget(top)

        h_cont = QHBoxLayout()
        left_scroll = QScrollArea(); left_scroll.setWidgetResizable(True); left_scroll.setFixedWidth(240)
        left = QWidget(); left_l = QVBoxLayout(left); left_l.setContentsMargins(5, 5, 5, 5)
        
        self.g_t = QGroupBox(self.tr["tools"]); v_t = QVBoxLayout(self.g_t); v_t.setSpacing(4); self.grp = QButtonGroup(self)
        self.btn_tools = []
        tools_data = [("tool_select", 0), ("tool_h_man", 7), ("tool_a_man", 8), ("tool_arr_s", 1), ("tool_arr_c", 2), ("tool_rect", 4), ("tool_circ", 5), ("tool_text", 3), ("tool_ball", 6)]
        for key, m in tools_data:
            b = QPushButton(self.tr[key]); b.setCheckable(True); self.grp.addButton(b); b.clicked.connect(lambda ch, m=m: self.set_tool_mode(m))
            v_t.addWidget(b); self.btn_tools.append((b, key))
            
        self.c_ball = QComboBox(); self.c_ball.addItems(self.tr["balls"]); self.c_ball.currentIndexChanged.connect(self.set_ball_type); v_t.addWidget(self.c_ball)
        self.grp.buttons()[0].setChecked(True); left_l.addWidget(self.g_t)

        self.g_e = QGroupBox(self.tr["edit_obj"]); self.g_e.setEnabled(False); self.v_e = QFormLayout(self.g_e); self.v_e.setVerticalSpacing(4)
        self.e_n = QLineEdit(); self.e_n.textChanged.connect(self.update_obj_name)
        self.lbl_col_target = QLabel(self.tr["color_target"]); self.color_target = QComboBox(); self.color_target.addItems([self.tr["col_target_jersey"], self.tr["col_target_text"]])
        color_layout = QHBoxLayout(); quick_colors = ["#ffffff", "#e74c3c", "#3498db", "#f1c40f", "#2ecc71", "#000000"]
        for qc in quick_colors:
            btn_c = QPushButton(); btn_c.setFixedSize(20, 20); btn_c.setStyleSheet(f"background-color: {qc}; border-radius: 10px; border: 1px solid #d4af37;")
            btn_c.clicked.connect(lambda checked, col=qc: self.apply_quick_color(col)); color_layout.addWidget(btn_c)
            
        self.b_c = QPushButton(self.tr["other"]); self.b_c.clicked.connect(self.change_obj_color); color_layout.addWidget(self.b_c)
        self.s_s = QSlider(Qt.Orientation.Horizontal); self.s_s.setRange(30, 250); self.s_s.setValue(100); self.s_s.valueChanged.connect(self.update_scale)
        self.s_r = QSlider(Qt.Orientation.Horizontal); self.s_r.setRange(0, 360); self.s_r.valueChanged.connect(self.update_rotation)
        self.s_c = QSlider(Qt.Orientation.Horizontal); self.s_c.setRange(-100, 100); self.s_c.valueChanged.connect(self.update_arrow_curve)
        m_l = QHBoxLayout(); self.b_m_h = QPushButton(self.tr["mir_h"]); self.b_m_h.clicked.connect(self.update_mirror_h)
        self.b_m_v = QPushButton(self.tr["mir_v"]); self.b_m_v.clicked.connect(self.update_mirror_v); m_l.addWidget(self.b_m_h); m_l.addWidget(self.b_m_v)
        
        self.btn_replace_img = QPushButton(self.tr["replace_img"]); self.btn_replace_img.setProperty("class", "ActionBtn"); self.btn_replace_img.clicked.connect(self.replace_player_image)
        self.btn_del = QPushButton(self.tr["delete_sel"]); self.btn_del.setProperty("class", "DangerBtn"); self.btn_del.clicked.connect(self.delete_active_item)

        self.lbl_name = QLabel(self.tr["name"]); self.v_e.addRow(self.lbl_name, self.e_n); self.v_e.addRow(self.lbl_col_target, self.color_target)
        self.lbl_col = QLabel(self.tr["color"]); self.v_e.addRow(self.lbl_col, color_layout); self.lbl_sz = QLabel(self.tr["size"]); self.v_e.addRow(self.lbl_sz, self.s_s)
        self.lbl_ang = QLabel(self.tr["angle"]); self.v_e.addRow(self.lbl_ang, self.s_r); self.lbl_mir = QLabel(self.tr["mirror"]); self.v_e.addRow(self.lbl_mir, m_l)
        self.lbl_crv = QLabel(self.tr["curve"]); self.v_e.addRow(self.lbl_crv, self.s_c); self.v_e.addRow("", self.btn_replace_img); self.v_e.addRow("", self.btn_del)
        left_l.addWidget(self.g_e); left_l.addStretch(); left_scroll.setWidget(left)
        h_cont.addWidget(left_scroll); h_cont.addWidget(self.board, 1)

        self.tabs = QTabWidget(); self.tabs.setFixedWidth(280)
        t_m = QWidget(); l_m = QVBoxLayout(t_m); l_m.setContentsMargins(5,5,5,5)
        self.m_l = QListWidget(); self.m_l.setViewMode(QListWidget.ViewMode.IconMode); self.m_l.setIconSize(QSize(60, 60)); self.m_l.itemClicked.connect(self.apply_memory_image)
        self.m_l.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu); self.m_l.customContextMenuRequested.connect(self.show_memory_menu)
        l_m.addWidget(self.m_l); self.b_u_img = QPushButton(self.tr["add_pc"]); self.b_u_img.setProperty("class", "ActionBtn"); self.b_u_img.clicked.connect(self.upload_to_memory); l_m.addWidget(self.b_u_img); self.tabs.addTab(t_m, self.tr["tab_img"])
        
        t_t = QWidget(); l_t = QVBoxLayout(t_t); l_t.setContentsMargins(5,5,5,5)
        self.t_l = QListWidget(); self.t_l.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu); self.t_l.customContextMenuRequested.connect(self.show_tactics_menu)
        l_t.addWidget(self.t_l); self.b_l_t = QPushButton(self.tr["load"]); self.b_l_t.clicked.connect(self.load_selected_tactic); l_t.addWidget(self.b_l_t); self.tabs.addTab(t_t, self.tr["tab_tac"])
        
        t_v = QWidget(); l_v = QVBoxLayout(t_v); l_v.setContentsMargins(5,5,5,5)
        self.v_l = QListWidget(); self.v_l.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu); self.v_l.customContextMenuRequested.connect(self.show_video_menu); self.v_l.itemDoubleClicked.connect(self.open_video)
        l_v.addWidget(self.v_l); self.b_v_u = QPushButton(self.tr["add_vid"]); self.b_v_u.setProperty("class", "ActionBtn"); self.b_v_u.clicked.connect(self.upload_video); l_v.addWidget(self.b_v_u); self.tabs.addTab(t_v, self.tr["tab_vid"])

        h_cont.addWidget(self.tabs); main_layout.addLayout(h_cont)
        
        QShortcut(QKeySequence("Ctrl+Z"), self).activated.connect(self.undo_action)
        QShortcut(QKeySequence("Ctrl+S"), self).activated.connect(self.save_tactics)
        
        self.load_memory_images(); self.load_tactics_list(); self.load_videos_list()
        self.push_state()
        self.showMaximized()

        QTimer.singleShot(1500, self.check_updates)
        QTimer.singleShot(100, self.fetch_remaining_days)

    def fetch_remaining_days(self):
        if not self.user_email:
            self.lbl_license.setText("Süre Yok")
            return
        try:
            url = f"https://firestore.googleapis.com/v1/projects/footyscopefg-df329/databases/(default)/documents/users/{self.user_email}"
            resp = requests.get(url, timeout=5)
            if resp.status_code == 200:
                db_data = resp.json()
                create_time_str = db_data.get("createTime", "")
                if not create_time_str:
                    self.lbl_license.setText("Hata")
                    return
                
                create_time_str = create_time_str.split(".")[0].replace("Z", "")
                create_time = datetime.strptime(create_time_str, "%Y-%m-%dT%H:%M:%S").replace(tzinfo=timezone.utc)
                server_date_str = resp.headers.get('Date')
                server_time = parsedate_to_datetime(server_date_str) if server_date_str else datetime.now(timezone.utc)
                
                days_passed = (server_time - create_time).days
                remaining = 180 - days_passed
                if remaining < 0: remaining = 0
                
                self.lbl_license.setText(f"⏳ Lisans: {remaining} Gün")
                if remaining <= 5:
                    self.lbl_license.setStyleSheet("background-color: #4a121d; color: #ff4d4d; border: 1px solid #ff4d4d; border-radius: 4px; padding: 4px; font-weight: bold; font-size: 11px;")
                else:
                    self.lbl_license.setStyleSheet("background-color: #2b0a11; color: #d4af37; border: 1px solid #b76e79; border-radius: 4px; padding: 4px; font-weight: bold; font-size: 11px;")
            else:
                self.lbl_license.setText("Hata")
        except Exception:
            self.lbl_license.setText("Bağlantı Yok")

    # Dosyanın en tepesinde:
APP_VERSION = "2.4"

# Sonra check_updates fonksiyonunun içinde:
def check_updates(self):
    current_version = APP_VERSION # Sabit rakam yerine yukarıdaki değişkeni kullan!
        
    try:
            url = "https://firestore.googleapis.com/v1/projects/footyscopefg-df329/databases/(default)/documents/system/updates"
            response = requests.get(url, timeout=5)
            
            if response.status_code == 200:
                data = response.json()
                fields = data.get('fields', {})
                
                if 'stringValue' in fields.get('latestVersion', {}):
                    latest_version = fields['latestVersion']['stringValue']
                else:
                    val = fields.get('latestVersion', {})
                    latest_version = str(val.get('doubleValue', val.get('integerValue', '1.0')))
                
                if sys.platform == "win32":
                    download_url = fields.get('downloadUrlWin', {}).get('stringValue', '')
                else:
                    download_url = fields.get('downloadUrlMac', {}).get('stringValue', '')
                
                QMessageBox.information(
                    self, 
                    "Sistem Röntgeni 🔍", 
                    f"Koddaki Sürüm: {current_version}\n"
                    f"Firebase Sürümü: {latest_version}\n"
                    f"İndirme Linki: {download_url}"
                )

                if current_version != latest_version and download_url:
                    reply = QMessageBox.question(
                        self, 
                        "Yeni Sürüm Bulundu 🚀", 
                        f"FootyscopeFG v{latest_version} hazır!\n\nŞu anki sürümünüz: v{current_version}\nUygulama otomatik güncellenip yeniden başlatılsın mı?",
                        QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
                    )
                    if reply == QMessageBox.StandardButton.Yes:
                        self.perform_autonomous_update(download_url)
            else:
                QMessageBox.warning(self, "Hata", f"Firebase HTTP Kodu: {response.status_code}")
                
    except Exception as e:
            QMessageBox.critical(self, "Çökme Yaşandı", f"Hata detayı:\n{str(e)}")

    def perform_autonomous_update(self, download_url):
        from PyQt6.QtWidgets import QProgressDialog
        from PyQt6.QtCore import Qt
        import tempfile
        import os
        import sys
        import subprocess
        import requests

        temp_dir = tempfile.gettempdir()
        is_win = sys.platform == "win32"
        temp_file_name = "FootyscopeFG_Update.exe" if is_win else "FootyscopeFG_Update.zip"
        temp_file_path = os.path.join(temp_dir, temp_file_name)

        try:
            # --- 1. ÇELİK YELEK: GOOGLE DRIVE LİNK DÜZELTİCİ ---
            if "drive.google.com" in download_url:
                file_id = ""
                if "id=" in download_url:
                    file_id = download_url.split("id=")[1].split("&")[0]
                elif "/file/d/" in download_url:
                    file_id = download_url.split("/file/d/")[1].split("/")[0]
                
                if file_id:
                    download_url = f"https://drive.google.com/uc?export=download&id={file_id}"

            session = requests.Session()
            response = session.get(download_url, stream=True, timeout=15)
            
            token = None
            for key, value in response.cookies.items():
                if key.startswith('download_warning'):
                    token = value
                    break
            
            if token:
                response = session.get(download_url, params={'confirm': token}, stream=True)

            total_size = int(response.headers.get('content-length', 0))

            progress = QProgressDialog("Yeni sürüm indiriliyor... Lütfen bekleyin.", "İptal", 0, total_size if total_size > 0 else 100, self)
            progress.setWindowTitle("FootyscopeFG Güncelleniyor")
            progress.setWindowModality(Qt.WindowModality.WindowModal)
            progress.setValue(0)
            progress.show()

            downloaded_size = 0
            with open(temp_file_path, "wb") as f:
                for chunk in response.iter_content(chunk_size=32768):
                    if progress.wasCanceled():
                        QMessageBox.warning(self, "İptal Edildi", "Güncelleme iptal edildi.")
                        return 
                    if chunk:
                        f.write(chunk)
                        downloaded_size += len(chunk)
                        if total_size > 0:
                            progress.setValue(downloaded_size)
                        else:
                            progress.setValue((progress.value() + 1) % 100)
                        QApplication.processEvents() 
            
            progress.setValue(progress.maximum())

            if getattr(sys, 'frozen', False):
                current_file = sys.executable 
                file_dir = os.path.dirname(current_file)
                file_name = os.path.basename(current_file)
                
                if is_win:
                    bat_path = os.path.join(temp_dir, "updater.bat")
                    bat_content = f"""@echo off
timeout /t 2 /nobreak > NUL
cd /d "{file_dir}"
del "{file_name}"
move /y "{temp_file_path}" "{file_name}"
start "" "{file_name}"
del "%~f0"
"""
                    with open(bat_path, "w", encoding="utf-8") as f:
                        f.write(bat_content)
                    subprocess.Popen(bat_path, shell=True)
                else: 
                    app_path = os.path.dirname(os.path.dirname(os.path.dirname(current_file)))
                    app_parent_dir = os.path.dirname(app_path)
                    sh_path = os.path.join(temp_dir, "updater.sh")
                    
                    # --- 2. ÇELİK YELEK: GÜVENLİ MAC GÜNCELLEME BETİĞİ ---
                    sh_content = f"""#!/bin/sh
sleep 2
# Inen dosya gercekten saglam bir ZIP mi kontrol et
if unzip -t "{temp_file_path}" > /dev/null 2>&1; then
    # Saglamsa eskiyi sil
    rm -rf "{app_path}"
    # Yenisini cikar
    unzip -o -q "{temp_file_path}" -d "{app_parent_dir}"
    rm -f "{temp_file_path}"
    # Mac'in Geliştirici Karantinasını kaldır (Uygulama direk açılsın diye)
    xattr -cr "{app_path}"
    # Uygulamayi ac
    open "{app_path}"
else
    # Dosya bozuksa eskiyi KESINLIKLE SILME!
    echo "Indirilen dosya bozuk veya Google Drive engeline takildi." > "{app_parent_dir}/GUNCELLEME_HATASI.txt"
fi
rm -- "$0"
"""
                    with open(sh_path, "w", encoding="utf-8") as f:
                        f.write(sh_content)
                    
                    # Script'e çalışma izni ver ve başlat
                    os.chmod(sh_path, 0o755)
                    subprocess.Popen(["sh", sh_path])
                
                sys.exit()
            else:
                QMessageBox.information(self, "Test Başarılı", f"Test modu: Dosya {temp_file_path} konumuna indi.")
                
        except Exception as e:
            QMessageBox.critical(self, "Hata", f"Güncelleme hatası:\n{str(e)}")

    def change_language(self, lang_name):
        self.tr = LANG.get(lang_name, LANG["Türkçe"])
        self.setWindowTitle(self.tr["title"])
        self.g1.setTitle(self.tr["pitch_group"])
        self.lbl_dir.setText(self.tr["dir"])
        self.lbl_surf.setText(self.tr["surface"])
        
        self.c_pitch.blockSignals(True)
        idx = self.c_pitch.currentIndex()
        self.c_pitch.clear()
        self.c_pitch.addItems(self.tr["pitch_types"])
        self.c_pitch.setCurrentIndex(idx)
        self.c_pitch.blockSignals(False)

        self.c_grass.blockSignals(True)
        idx2 = self.c_grass.currentIndex()
        self.c_grass.clear()
        self.c_grass.addItems(self.tr["grass_types"])
        self.c_grass.setCurrentIndex(idx2)
        self.c_grass.blockSignals(False)

        self.g_h.setTitle(self.tr["home"])
        self.g_a.setTitle(self.tr["away"])
        self.c_f_h.setItemText(0, self.tr["select_tactic"])
        self.c_f_a.setItemText(0, self.tr["select_tactic"])
        self.b_h1.setText(self.tr["color1"])
        self.b_h2.setText(self.tr["color2"])
        self.b_a1.setText(self.tr["color1"])
        self.b_a2.setText(self.tr["color2"])
        self.b_a_h.setText(self.tr["apply"])
        self.b_a_a.setText(self.tr["apply"])
        
        self.g_o.setTitle(self.tr["management"])
        self.b_s.setText(self.tr["save"])
        self.b_p.setText(self.tr["pptx"])
        self.b_pdf.setText(self.tr["pdf"])
        self.b_logout.setText(self.tr["logout"])
        if not self.is_recording:
            self.b_v.setText(self.tr["record"])
        else:
            self.b_v.setText(self.tr["stop_record"])

        self.g_t.setTitle(self.tr["tools"])
        for btn, key in self.btn_tools:
            btn.setText(self.tr[key])
            
        self.c_ball.blockSignals(True)
        idx3 = self.c_ball.currentIndex()
        self.c_ball.clear()
        self.c_ball.addItems(self.tr["balls"])
        self.c_ball.setCurrentIndex(idx3)
        self.c_ball.blockSignals(False)

        self.g_e.setTitle(self.tr["edit_obj"])
        self.lbl_col_target.setText(self.tr["color_target"])
        self.color_target.blockSignals(True)
        self.color_target.setItemText(0, self.tr["col_target_jersey"])
        self.color_target.setItemText(1, self.tr["col_target_text"])
        self.color_target.blockSignals(False)

        self.lbl_name.setText(self.tr["name"])
        self.lbl_col.setText(self.tr["color"])
        self.lbl_sz.setText(self.tr["size"])
        self.lbl_ang.setText(self.tr["angle"])
        self.lbl_mir.setText(self.tr["mirror"])
        self.b_m_h.setText(self.tr["mir_h"])
        self.b_m_v.setText(self.tr["mir_v"])
        self.lbl_crv.setText(self.tr["curve"])
        self.b_c.setText(self.tr["other"])
        self.btn_replace_img.setText(self.tr["replace_img"])
        self.btn_del.setText(self.tr["delete_sel"])

        self.tabs.setTabText(0, self.tr["tab_img"])
        self.tabs.setTabText(1, self.tr["tab_tac"])
        self.tabs.setTabText(2, self.tr["tab_vid"])
        self.b_u_img.setText(self.tr["add_pc"])
        self.b_l_t.setText(self.tr["load"])
        self.b_v_u.setText(self.tr["add_vid"])

    def refresh_pitch(self):
        old_w, old_h = self.board.p_w, self.board.p_h
        rel_pos = []
        for item in self.board.scene.items():
            if isinstance(item, (Player, Football, EditableText)):
                rel_pos.append((item, item.x()/old_w, item.y()/old_h))
            elif isinstance(item, Arrow):
                rel_pos.append((item, item.start_p.x()/old_w, item.start_p.y()/old_h, item.end_p.x()/old_w, item.end_p.y()/old_h))

        self.board.draw_pitch(self.c_pitch.currentIndex(), self.c_grass.currentIndex())
        new_w, new_h = self.board.p_w, self.board.p_h

        for data in rel_pos:
            item = data[0]
            if isinstance(item, (Player, Football, EditableText)):
                item.setPos(data[1]*new_w, data[2]*new_h)
            elif isinstance(item, Arrow):
                item.update_arrow(QPointF(data[1]*new_w, data[2]*new_h), QPointF(data[3]*new_w, data[4]*new_h))

        self.board.fitInView(self.board.scene.sceneRect(), Qt.AspectRatioMode.KeepAspectRatio)

    def get_scene_state(self):
        d = []
        for it in self.board.scene.items():
            if isinstance(it, Player):
                d.append({"t": "p", "x": it.x(), "y": it.y(), "n": it.name_tag.toPlainText(), "c1": it.color1.name(), "c2": it.color2.name() if it.color2 else None, "gk": it.is_gk, "img": getattr(it, 'image_path', PLAYER_ASSET), "s": it.scale(), "r": it.rotation(), "mh": it.is_mirrored_h, "mv": it.is_mirrored_v, "team": getattr(it, 'team', 'none'), "isc": getattr(it, 'is_custom', False), "tc": it.name_tag.defaultTextColor().name()})
            elif isinstance(it, Football):
                d.append({"t": "b", "x": it.x(), "y": it.y(), "type": getattr(it, 'b_idx', 0), "s": it.scale()})
            elif isinstance(it, Arrow):
                d.append({"t": "a", "spx": it.start_p.x(), "spy": it.start_p.y(), "epx": it.end_p.x(), "epy": it.end_p.y(), "curv": it.is_curved, "coff": getattr(it, 'curve_offset', 0.3), "col": it.color.name()})
            elif isinstance(it, EditableText):
                d.append({"t": "text", "x": it.x(), "y": it.y(), "content": it.toPlainText(), "color": it.defaultTextColor().name(), "s": it.scale(), "r": it.rotation()})
            elif isinstance(it, CustomRect):
                d.append({"t": "rect", "x": it.rect().x(), "y": it.rect().y(), "w": it.rect().width(), "h": it.rect().height(), "pos_x": it.x(), "pos_y": it.y(), "col": it.brush().color().name() if it.brush().style() != Qt.BrushStyle.NoBrush else "#ffffff"})
            elif isinstance(it, CustomEllipse):
                d.append({"t": "ell", "x": it.rect().x(), "y": it.rect().y(), "w": it.rect().width(), "h": it.rect().height(), "pos_x": it.x(), "pos_y": it.y(), "col": it.brush().color().name() if it.brush().style() != Qt.BrushStyle.NoBrush else "#ffffff"})
        return d

    def load_scene_state(self, data):
        for i in self.board.scene.items():
            if isinstance(i, (Player, Football, Arrow, EditableText, CustomRect, CustomEllipse)):
                self.board.scene.removeItem(i)
        for d in data:
            if d["t"] == "p":
                p = Player(d.get("img", GK_ASSET if d["gk"] else PLAYER_ASSET), d["n"], d["c1"], d.get("c2"), d["x"], d["y"], d["gk"], is_custom=d.get("isc", False))
                p.team = d.get("team", 'none')
                p.setScale(d.get("s", 1.0))
                p.setRotation(d.get("r", 0.0))
                p.set_mirrored(d.get("mh", False), d.get("mv", False))
                p.name_tag.setDefaultTextColor(QColor(d.get("tc", "#ffffff")))
                self.board.scene.addItem(p)
            elif d["t"] == "b":
                b = Football(d["x"], d["y"], d.get("type", 0))
                b.setScale(d.get("s", 1.0))
                self.board.scene.addItem(b)
            elif d["t"] == "a":
                a = Arrow(QPointF(d["spx"], d["spy"]), QPointF(d["epx"], d["epy"]), d["curv"], d["col"])
                a.curve_offset = d.get("coff", 0.3)
                a.update_arrow(a.start_p, a.end_p)
                self.board.scene.addItem(a)
            elif d["t"] == "text":
                t = EditableText(d["content"], d["x"], d["y"])
                t.setDefaultTextColor(QColor(d["color"]))
                t.setScale(d.get("s", 1.0))
                t.setRotation(d.get("r", 0.0))
                self.board.scene.addItem(t)
            elif d["t"] == "rect":
                r = CustomRect(QRectF(d["x"], d["y"], d["w"], d["h"]))
                r.setPos(d["pos_x"], d["pos_y"])
                r.setBrush(QBrush(QColor(d["col"])))
                self.board.scene.addItem(r)
            elif d["t"] == "ell":
                e = CustomEllipse(QRectF(d["x"], d["y"], d["w"], d["h"]))
                e.setPos(d["pos_x"], d["pos_y"])
                e.setBrush(QBrush(QColor(d["col"])))
                self.board.scene.addItem(e)

    def push_state(self):
        state = self.get_scene_state()
        if self.history_index < len(self.history) - 1:
            self.history = self.history[:self.history_index + 1]
        self.history.append(state)
        self.history_index += 1

    def undo_action(self):
        if self.history_index > 0:
            self.history_index -= 1
            self.load_scene_state(self.history[self.history_index])

    def load_videos_list(self):
        self.v_l.clear()
        for f in os.listdir(VIDEO_DIR):
            if f.lower().endswith(('.mp4', '.mov', '.avi', '.mkv')):
                item = QListWidgetItem(f"🎬 {f}")
                item.setData(Qt.ItemDataRole.UserRole, os.path.join(VIDEO_DIR, f))
                self.v_l.addItem(item)
                
    def upload_video(self):
        ps, _ = QFileDialog.getOpenFileNames(self, "Video Seç", "", "Videolar (*.mp4 *.mov *.avi *.mkv)")
        for p in ps:
            shutil.copy2(p, VIDEO_DIR)
        self.load_videos_list()
        
    def open_video(self, item):
        path = item.data(Qt.ItemDataRole.UserRole)
        if sys.platform == "win32":
            os.startfile(path)
        else:
            subprocess.call(["open", path])
            
    def show_video_menu(self, pos):
        item = self.v_l.itemAt(pos)
        if item:
            menu = QMenu()
            open_act = menu.addAction(self.tr["ctx_play"])
            ren_act = menu.addAction(self.tr["ctx_ren"])
            del_act = menu.addAction(self.tr["ctx_del_vid"])
            action = menu.exec(self.v_l.mapToGlobal(pos))
            path = item.data(Qt.ItemDataRole.UserRole)
            if action == open_act:
                self.open_video(item)
            elif action == ren_act:
                new_name, ok = QInputDialog.getText(self, self.tr["ctx_ren"], "İsim:")
                if ok and new_name:
                    ext = os.path.splitext(path)[1]
                    new_path = os.path.join(VIDEO_DIR, f"{new_name}{ext}")
                    os.rename(path, new_path)
                    self.load_videos_list()
            elif action == del_act:
                os.remove(path)
                self.load_videos_list()

    def load_team_logo(self, team):
        path, _ = QFileDialog.getOpenFileName(self, "Logo", "", "Resimler (*.png *.jpg *.jpeg)")
        if path:
            pixmap = QPixmap(path).scaled(24, 24, Qt.AspectRatioMode.KeepAspectRatio, Qt.TransformationMode.SmoothTransformation)
            if team == "home":
                self.lbl_logo_h.setPixmap(pixmap)
            else:
                self.lbl_logo_a.setPixmap(pixmap)

    def show_memory_menu(self, pos):
        item = self.m_l.itemAt(pos)
        if item:
            menu = QMenu()
            del_action = menu.addAction(self.tr["ctx_del_pc"])
            action = menu.exec(self.m_l.mapToGlobal(pos))
            if action == del_action:
                path = item.data(Qt.ItemDataRole.UserRole)
                try:
                    os.remove(path)
                except:
                    pass
                self.load_memory_images()

    def show_tactics_menu(self, pos):
        item = self.t_l.itemAt(pos)
        if item:
            menu = QMenu()
            del_action = menu.addAction(self.tr["ctx_del_tac"])
            action = menu.exec(self.t_l.mapToGlobal(pos))
            if action == del_action:
                path = os.path.join(TACTICS_DIR, f"{item.text()}.json")
                try:
                    os.remove(path)
                except:
                    pass
                self.load_tactics_list()

    def replace_player_image(self):
        if isinstance(self.active_item, Player):
            path, _ = QFileDialog.getOpenFileName(self, "Görsel", "", "Resimler (*.png *.jpg *.jpeg)")
            if path: 
                self.active_item.is_custom = True
                self.active_item.load_image(path)
                self.push_state()

    def logout(self):
        if os.path.exists(AUTH_FILE):
            os.remove(AUTH_FILE)
        QMessageBox.information(self, "Çıkış", "Oturumunuz kapatıldı. Uygulama yeniden başlatılıyor...")
        QApplication.quit()

    def delete_active_item(self):
        if self.active_item: 
            self.board.scene.removeItem(self.active_item)
            self.active_item = None
            self.board.scene.clearSelection()
            self.g_e.setEnabled(False)
            self.push_state()

    def set_btn_color(self, b, h): 
        text_c = "black" if h.lower() in ["#ffffff", "#fff", "white"] else "white"
        b.setStyleSheet(f"background-color: {h}; color: {text_c}; border-radius: 4px; border: 1px solid #d4af37;")
        b.setProperty("hex", h)
        
    def pick_team_color(self, b):
        c = QColorDialog.getColor()
        if c.isValid():
            self.set_btn_color(b, c.name())

    def set_ball_type(self, idx):
        self.current_ball_type = idx
        
    def set_tool_mode(self, m):
        if m == 6:
            self.board.scene.addItem(Football(self.board.p_w/2, self.board.p_h/2, self.current_ball_type))
            self.grp.buttons()[0].setChecked(True)
            self.board.current_mode = 0
            self.push_state()
        else:
            self.board.current_mode = m
            self.board.scene.clearSelection()
    
    def on_selection_changed(self):
        items = self.board.scene.selectedItems()
        if items:
            self.active_item = items[0]
            if isinstance(self.active_item, QGraphicsTextItem) and isinstance(self.active_item.parentItem(), Player):
                self.active_item = self.active_item.parentItem()
            self.g_e.setEnabled(True)
            self.s_c.setVisible(isinstance(self.active_item, Arrow) and self.active_item.is_curved)
            self.btn_replace_img.setVisible(isinstance(self.active_item, Player))
            self.lbl_col_target.setVisible(isinstance(self.active_item, Player))
            self.color_target.setVisible(isinstance(self.active_item, Player))
            self.e_n.setText(self.active_item.name_tag.toPlainText() if isinstance(self.active_item, Player) else "")
            
            if isinstance(self.active_item, (Player, Football, CustomRect, CustomEllipse, EditableText)):
                self.s_s.blockSignals(True)
                self.s_s.setValue(int(self.active_item.scale() * 100))
                self.s_s.blockSignals(False)
            if isinstance(self.active_item, (Player, EditableText)):
                self.s_r.blockSignals(True)
                self.s_r.setValue(int(self.active_item.rotation()))
                self.s_r.blockSignals(False)
        else:
            self.g_e.setEnabled(False)

    def apply_formation(self, team):
        is_h = (team == "home")
        c = self.c_f_h if is_h else self.c_f_a
        f_idx = c.currentIndex()
        if f_idx == 0: return
        f = c.itemText(f_idx)
        
        for item in self.board.scene.items():
            if isinstance(item, Player) and getattr(item, 'team', None) == team:
                self.board.scene.removeItem(item)
                
        is_v = self.c_pitch.currentIndex() in [1, 3] 
        w, h = self.board.p_w, self.board.p_h
        
        def p_p(px, py):
            if not is_v:
                if is_h: return (px * w, py * h)
                else: return ((1 - px) * w, (1 - py) * h)
            else:
                if is_h: return (py * w, (1 - px) * h)
                else: return ((1 - py) * w, px * h)

        c1 = self.b_h1.property("hex") if is_h else self.b_a1.property("hex")
        c2 = self.b_h2.property("hex") if is_h else self.b_a2.property("hex")
        ps = []
        if f == "4-4-2": ps = [(0.1, 0.5, True), (0.3, 0.2, False), (0.3, 0.4, False), (0.3, 0.6, False), (0.3, 0.8, False), (0.6, 0.2, False), (0.6, 0.4, False), (0.6, 0.6, False), (0.6, 0.8, False), (0.85, 0.4, False), (0.85, 0.6, False)]
        elif f == "4-3-3": ps = [(0.1, 0.5, True), (0.3, 0.2, False), (0.3, 0.4, False), (0.3, 0.6, False), (0.3, 0.8, False), (0.5, 0.5, False), (0.6, 0.3, False), (0.6, 0.7, False), (0.8, 0.2, False), (0.85, 0.5, False), (0.8, 0.8, False)]
        elif f == "3-5-2": ps = [(0.1, 0.5, True), (0.3, 0.3, False), (0.3, 0.5, False), (0.3, 0.7, False), (0.55, 0.15, False), (0.5, 0.35, False), (0.5, 0.65, False), (0.55, 0.85, False), (0.65, 0.5, False), (0.85, 0.4, False), (0.85, 0.6, False)]
        elif f == "4-2-3-1": ps = [(0.1, 0.5, True), (0.3, 0.2, False), (0.3, 0.4, False), (0.3, 0.6, False), (0.3, 0.8, False), (0.5, 0.35, False), (0.5, 0.65, False), (0.7, 0.2, False), (0.7, 0.5, False), (0.7, 0.8, False), (0.85, 0.5, False)]
        elif f == "3-4-1-2": ps = [(0.1, 0.5, True), (0.3, 0.3, False), (0.3, 0.5, False), (0.3, 0.7, False), (0.5, 0.2, False), (0.5, 0.4, False), (0.5, 0.6, False), (0.5, 0.8, False), (0.7, 0.5, False), (0.85, 0.4, False), (0.85, 0.6, False)]
        elif f == "4-1-4-1": ps = [(0.1, 0.5, True), (0.3, 0.2, False), (0.3, 0.4, False), (0.3, 0.6, False), (0.3, 0.8, False), (0.45, 0.5, False), (0.65, 0.2, False), (0.65, 0.4, False), (0.65, 0.6, False), (0.65, 0.8, False), (0.85, 0.5, False)]

        pre = self.tr["ev_pre"] + " " if is_h else self.tr["dep_pre"] + " "
        for i, (px, py, gk) in enumerate(ps):
            x, y = p_p(px, py)
            p = Player(GK_ASSET if gk else PLAYER_ASSET, f"{pre}{i}", c1, c2, x, y, gk)
            p.team = team
            if not is_v:
                p.set_mirrored(not is_h, False)
            else:
                p.setRotation(270 if is_h else 90)
            self.board.scene.addItem(p)
        self.push_state()

    def update_obj_name(self, t): 
        if isinstance(self.active_item, Player):
            self.active_item.name_tag.setPlainText(t)
            self.active_item.update_name_pos()
            
    def update_arrow_curve(self, v): 
        if isinstance(self.active_item, Arrow) and self.active_item.is_curved:
            self.active_item.curve_offset = v/100.0
            self.active_item.update_arrow(self.active_item.start_p, self.active_item.end_p)
            
    def update_scale(self, v): 
        if self.active_item: self.active_item.setScale(v/100.0)
        
    def update_rotation(self, v): 
        if self.active_item: self.active_item.setRotation(v)
    
    def update_mirror_h(self):
        if isinstance(self.active_item, Player):
            self.active_item.set_mirrored(not self.active_item.is_mirrored_h, self.active_item.is_mirrored_v)
        elif isinstance(self.active_item, Arrow):
            self.active_item.curve_offset *= -1
            self.active_item.update_arrow(self.active_item.start_p, self.active_item.end_p)

    def update_mirror_v(self):
        if isinstance(self.active_item, Player):
            self.active_item.set_mirrored(self.active_item.is_mirrored_h, not self.active_item.is_mirrored_v)
        elif isinstance(self.active_item, Arrow):
            sp, ep = self.active_item.start_p, self.active_item.end_p
            self.active_item.update_arrow(ep, sp)

    def apply_quick_color(self, hex_color):
        if not self.active_item: return
        self._apply_color_to_item(QColor(hex_color))

    def change_obj_color(self):
        if not self.active_item: return
        c = QColorDialog.getColor()
        if c.isValid():
            self._apply_color_to_item(c)

    def _apply_color_to_item(self, c):
        if isinstance(self.active_item, EditableText): 
            self.active_item.setDefaultTextColor(c)
        elif isinstance(self.active_item, Player): 
            if self.color_target.currentIndex() == 1: 
                self.active_item.name_tag.setDefaultTextColor(c)
            else: 
                self.active_item.color1 = c
                self.active_item.color2 = None
                self.active_item.load_image(self.active_item.image_path)
        elif isinstance(self.active_item, Arrow): 
            self.active_item.color = c
            self.active_item.update_arrow(self.active_item.start_p, self.active_item.end_p)
        elif isinstance(self.active_item, (CustomRect, CustomEllipse, Football)): 
            self.active_item.setBrush(QBrush(c))
        self.push_state()

    def load_memory_images(self):
        self.m_l.clear()
        for f in os.listdir(MEMORY_DIR):
            if f.lower().endswith(('.png', '.jpg', '.jpeg')):
                item = QListWidgetItem(QIcon(os.path.join(MEMORY_DIR, f)), "")
                item.setData(Qt.ItemDataRole.UserRole, os.path.join(MEMORY_DIR, f))
                self.m_l.addItem(item)
    
    def upload_to_memory(self):
        ps, _ = QFileDialog.getOpenFileNames(self, "Görsel Seç", "", "Resimler (*.png *.jpg *.jpeg)")
        for p in ps: shutil.copy2(p, MEMORY_DIR)
        self.load_memory_images()

    def apply_memory_image(self, it):
        img_path = it.data(Qt.ItemDataRole.UserRole)
        if isinstance(self.active_item, Player): 
            self.active_item.is_custom = True
            self.active_item.load_image(img_path)
            self.push_state()
        else:
            new_img = Player(img_path, "Yeni Görsel", "#ffffff", None, self.board.p_w/2, self.board.p_h/2, is_custom=True)
            self.board.scene.addItem(new_img)
            self.board.scene.clearSelection()
            new_img.setSelected(True)
            self.push_state()

    def load_tactics_list(self):
        self.t_l.clear()
        for f in os.listdir(TACTICS_DIR):
            if f.endswith('.json'):
                self.t_l.addItem(f.replace('.json', ''))
    
    def save_tactics(self):
        n, ok = QInputDialog.getText(self, "Kaydet", "Taktik Adı:")
        if ok and n:
            d = self.get_scene_state()
            with open(os.path.join(TACTICS_DIR, f"{n}.json"), "w") as f:
                json.dump(d, f)
            self.load_tactics_list()
            
    def load_selected_tactic(self):
        it = self.t_l.currentItem()
        if not it: return
        with open(os.path.join(TACTICS_DIR, f"{it.text()}.json"), "r") as f:
            data = json.load(f)
        self.load_scene_state(data)
        self.push_state()

    def export_to_pptx(self):
        try:
            prs = Presentation()
            scene_rect = self.board.scene.sceneRect()
            aspect_ratio = scene_rect.width() / scene_rect.height()
            h_inch = 7.5
            w_inch = 7.5 * aspect_ratio
            prs.slide_width = Inches(w_inch)
            prs.slide_height = Inches(h_inch)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            all_items = self.board.scene.items()
            dynamic_items = [i for i in all_items if isinstance(i, (Player, EditableText, Arrow, Football, CustomRect, CustomEllipse))]
            for i in dynamic_items: i.hide()
            
            bg_path = os.path.join(APP_DIR, "temp_bg.png")
            img = QImage(int(scene_rect.width()), int(scene_rect.height()), QImage.Format.Format_ARGB32)
            img.fill(Qt.GlobalColor.transparent)
            p = QPainter(img)
            self.board.scene.render(p, QRectF(0, 0, scene_rect.width(), scene_rect.height()), scene_rect)
            p.end()
            img.save(bg_path)
            
            slide.shapes.add_picture(bg_path, 0, 0, width=Inches(w_inch), height=Inches(h_inch))
            
            for i in all_items: i.hide()
            old_bg = self.board.scene.backgroundBrush()
            self.board.scene.setBackgroundBrush(Qt.GlobalColor.transparent)
            scale_x, scale_y = w_inch / scene_rect.width(), h_inch / scene_rect.height()

            for item in dynamic_items:
                item.show()
                rect = item.sceneBoundingRect()
                if rect.width() > 0 and rect.height() > 0:
                    rect.adjust(-5, -5, 5, 5)
                    i_p = os.path.join(APP_DIR, f"temp_item_{id(item)}.png")
                    i_i = QImage(int(rect.width()), int(rect.height()), QImage.Format.Format_ARGB32)
                    i_i.fill(Qt.GlobalColor.transparent)
                    p2 = QPainter(i_i)
                    self.board.scene.render(p2, QRectF(0, 0, rect.width(), rect.height()), rect)
                    p2.end()
                    i_i.save(i_p)
                    left = Inches((rect.x() - scene_rect.x()) * scale_x)
                    top = Inches((rect.y() - scene_rect.y()) * scale_y)
                    slide.shapes.add_picture(i_p, left, top, width=Inches(rect.width() * scale_x), height=Inches(rect.height() * scale_y))
                    os.remove(i_p)
                item.hide()
                
            for i in all_items: i.show()
            self.board.scene.setBackgroundBrush(old_bg)
            os.remove(bg_path)
            
            p_out = os.path.join(os.path.expanduser("~"), "Desktop", "Hareketli_Taktik.pptx")
            prs.save(p_out)
            QMessageBox.information(self, "OK", "PPTX Masaüstüne kaydedildi.")
        except Exception as e:
            for i in self.board.scene.items(): i.show()
            QMessageBox.critical(self, "Hata", str(e))

    def export_to_pdf(self):
        try:
            path, _ = QFileDialog.getSaveFileName(self, "PDF", os.path.join(os.path.expanduser("~"), "Desktop", "Taktik_Raporu.pdf"), "PDF (*.pdf)")
            if not path: return
            writer = QPdfWriter(path)
            writer.setPageSize(QPageSize(QPageSize.PageSizeId.A4))
            writer.setPageOrientation(QPageLayout.Orientation.Landscape)
            painter = QPainter(writer)
            self.board.scene.render(painter, QRectF(), self.board.scene.sceneRect())
            painter.end()
            QMessageBox.information(self, "OK", "PDF Masaüstüne kaydedildi.")
        except Exception as e:
            QMessageBox.critical(self, "Hata", str(e))

    def toggle_recording(self):
        if not self.is_recording:
            self.board.scene.clearSelection()
            self.is_recording = True
            self.b_v.setText(self.tr["stop_record"])
            self.b_v.setStyleSheet("background: #f85149; color: white")
            size = self.board.viewport().size()
            p = os.path.join(os.path.expanduser("~"), "Desktop", "Animasyon.mp4")
            self.video_writer = cv2.VideoWriter(p, cv2.VideoWriter_fourcc(*'mp4v'), 20, (size.width(), size.height()))
            self.record_timer.start(50)
        else:
            self.is_recording = False
            self.record_timer.stop()
            self.video_writer.release()
            self.b_v.setText(self.tr["record"])
            self.b_v.setStyleSheet("")
            QMessageBox.information(self, "OK", "Video Masaüstüne kaydedildi.")
    
    def record_frame(self):
        img = QImage(self.board.viewport().size(), QImage.Format.Format_RGB32)
        p = QPainter(img)
        self.board.viewport().render(p)
        p.end()
        ptr = img.bits()
        ptr.setsize(img.sizeInBytes())
        frame = np.frombuffer(ptr, np.uint8).reshape((img.height(), img.width(), 4))
        self.video_writer.write(cv2.cvtColor(frame, cv2.COLOR_BGRA2BGR))

class LoginDialog(QDialog):
    def __init__(self):
        super().__init__()
        self.user_email = ""
        self.setWindowTitle("FootyScopeFG - Lisans Girişi")
        self.setFixedSize(350, 240)
        self.setStyleSheet(LUXURY_THEME_QSS)
        
        layout = QVBoxLayout(self)
        self.lbl_info = QLabel("Lütfen web sitesinden aldığınız\nlisans anahtarını girin.", self)
        self.lbl_info.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.email_input = QLineEdit(self)
        self.email_input.setPlaceholderText("E-posta Adresiniz")
        self.key_input = QLineEdit(self)
        self.key_input.setPlaceholderText("Lisans Anahtarı (Örn: FS-1234)")
        self.remember_cb = QCheckBox("Oturumu Açık Tut (Beni Hatırla)", self)
        self.remember_cb.setChecked(True)
        self.btn_login = QPushButton("Giriş Yap", self)
        self.btn_login.setProperty("class", "ActionBtn")
        self.btn_login.clicked.connect(self.check_license)
        
        layout.addWidget(self.lbl_info)
        layout.addWidget(self.email_input)
        layout.addWidget(self.key_input)
        layout.addWidget(self.remember_cb)
        layout.addWidget(self.btn_login)

    def check_license(self):
        email_text = self.email_input.text().strip().lower()
        key_text = self.key_input.text().strip()
        if email_text == "" or key_text == "":
            QMessageBox.warning(self, "Uyarı", "E-posta ve Lisans Anahtarı boş bırakılamaz!")
            return
            
        url = f"https://firestore.googleapis.com/v1/projects/footyscopefg-df329/databases/(default)/documents/users/{email_text}"
        try:
            response = requests.get(url, timeout=5)
            if response.status_code == 200:
                data = response.json()
                if 'fields' in data and 'licenseKey' in data['fields']:
                    if key_text == data['fields']['licenseKey']['stringValue']:
                        
                        create_time_str = data.get("createTime", "")
                        if create_time_str:
                            create_time_str = create_time_str.split(".")[0].replace("Z", "")
                            create_time = datetime.strptime(create_time_str, "%Y-%m-%dT%H:%M:%S").replace(tzinfo=timezone.utc)
                            server_date_str = response.headers.get('Date')
                            server_time = parsedate_to_datetime(server_date_str) if server_date_str else datetime.now(timezone.utc)
                     
                            if (server_time - create_time).days >= 180:
                                QMessageBox.critical(self, "Lisans Süresi Doldu 🛑", "Lisans Süreniz Dolmuştur, Lütfen Yöneticiyle İletişime Geçin")
                                sys.exit()
                                return
                            
                        if self.remember_cb.isChecked():
                            try:
                                with open(AUTH_FILE, "w") as f:
                                    json.dump({"email": email_text, "key": key_text}, f)
                            except:
                                pass
                        
                        self.user_email = email_text
                        self.accept()
                        return
            QMessageBox.critical(self, "Hata", "E-posta veya Lisans Anahtarı geçersiz!")
        except requests.exceptions.RequestException:
            QMessageBox.critical(self, "Bağlantı Hatası", "Sunucuya ulaşılamıyor.")

def attempt_auto_login():
    if os.path.exists(AUTH_FILE):
        try:
            with open(AUTH_FILE, "r") as f:
                data = json.load(f)
            saved_email = data.get("email", "")
            saved_key = data.get("key", "")
            if saved_email and saved_key:
                url = f"https://firestore.googleapis.com/v1/projects/footyscopefg-df329/databases/(default)/documents/users/{saved_email}"
                resp = requests.get(url, timeout=3)
                if resp.status_code == 200:
                    db_data = resp.json()
                    if 'fields' in db_data and 'licenseKey' in db_data['fields'] and saved_key == db_data['fields']['licenseKey']['stringValue']:
                        
                        create_time_str = db_data.get("createTime", "")
                        if create_time_str:
                            create_time_str = create_time_str.split(".")[0].replace("Z", "")
                            create_time = datetime.strptime(create_time_str, "%Y-%m-%dT%H:%M:%S").replace(tzinfo=timezone.utc)
                            server_date_str = resp.headers.get('Date')
                            server_time = parsedate_to_datetime(server_date_str) if server_date_str else datetime.now(timezone.utc)
                            
                            if (server_time - create_time).days >= 180:
                                return None
                        return saved_email
        except:
            pass
    return None

if __name__ == "__main__":
    app = QApplication(sys.argv)
    
    auto_email = attempt_auto_login()
    if auto_email:
        window = MainWindow(auto_email)
        window.show()
        sys.exit(app.exec())
    else:
        login = LoginDialog()
        if login.exec() == QDialog.DialogCode.Accepted:
            window = MainWindow(login.user_email)
            window.show()
            sys.exit(app.exec())